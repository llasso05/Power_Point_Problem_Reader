import os
from pptx import Presentation
from openpyxl import Workbook
from PIL import Image
import pytesseract
import io


# function
def ppt2xls(ppt_file, slide_numbers):
    presentation = Presentation(ppt_file)

    # creating workbook with pyxl
    workbook = Workbook()

    # Iterate the slide numbers, input must be iterable
    for i, slide_number in enumerate(slide_numbers, start=1):
        # dealing with index 0 vs ppt index
        slide = presentation.slides[slide_number-1]

        # new sheet
        sheet_name = f'Slide {slide_number}'
        if i == 1:
            sheet = workbook.active
            sheet.title = sheet_name
        else:
            sheet = workbook.create_sheet(title=sheet_name)
        
        # Getting the text
        text = ''
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

        # Writing text on A1
        sheet['A1'] = text.strip()

        # check for images
        image_index = 1
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = io.BytesIO(image.blob)

                image_name = f"{os.path.splitext(ppt_file)[0]}_slide_{slide_number}_img_{image_index}.png"
                with open(image_name, "wb") as f:
                    f.write(image_bytes.getvalue())

                # using OCR pytesseract to extract table from pictures
                extracted_text = pytesseract.image_to_data(Image.open(image_name), output_type=pytesseract.Output.DICT)
                # Process the extracted data to write into the Excel sheet
                for j in range(len(extracted_text['text'])):
                    if int(extracted_text['conf'][j]) > 60:  # Confidence threshold
                        x, y, w, h = extracted_text['left'][j], extracted_text['top'][j], extracted_text['width'][j], extracted_text['height'][j]
                        text = extracted_text['text'][j]
                        if text.strip():
                            # Example of writing to specific cells, modify as needed
                            sheet.append([text])  # Write extracted text in the next available row
                
                image_index +=1

    # Creating final xls file
    excel_file = os.path.splitext(ppt_file)[0] + ".xlsx"
    workbook.save(excel_file)
    print(f"excel file saved as: {excel_file}")

# test 
ppt_file = 'Clase NO. 8_ Regresión Logística.pptx'
slide_numbers = [18,19]
ppt2xls(ppt_file,slide_numbers)